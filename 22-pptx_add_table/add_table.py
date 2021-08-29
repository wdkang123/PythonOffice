from pptx import Presentation
from pptx.util import Inches


# 向PPT中插入表格
ppt = Presentation()
layout = ppt.slide_layouts[6]
slide = ppt.slides.add_slide(layout)

rows = 2
cols = 2
left = Inches(3.5)
top = Inches(4.5)
width = Inches(6)
height = Inches(0.8)

# 添加表格 获取表格类
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# 第一列宽度
table.columns[0].width = Inches(2.0)
# 第二列宽度
table.columns[1].width = Inches(4.0)

table.cell(0, 0).text = "第一行第一列"
table.cell(0, 1).text = "第一行第二列"
table.cell(1, 0).text = "第二行第一列"

ppt.save("add_table.pptx")
