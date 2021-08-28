from pptx import Presentation


# PPT 文件对象
ppt = Presentation()
# 遍历所有的布局
for layout in ppt.slide_layouts:
  # 为该PPT文件添加使用某种布局的幻灯片
  slide = ppt.slides.add_slide(layout)

# 保存ppt
ppt.save("show_all_layout.pptx")
