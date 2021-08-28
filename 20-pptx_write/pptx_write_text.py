from pptx import Presentation


ppt = Presentation()

# 幻灯片布局 选择第一种默认布局
slide_layout = ppt.slide_layouts[0]

# slide对象为一页幻灯片 一个ppt文件中可以有多页幻灯片
slide = ppt.slides.add_slide(slide_layout)
# 取本页幻灯片的title占位符
title = slide.shapes.title
# 向title文本框中插入文字
title.text = "我是标题"
# 取出本页幻灯片 第二个文本框
subtitle = slide.placeholders[1]
# 向第二个文本框插入文字
subtitle.text = "正文框"

# 第二页ppt 采用不同的布局
slide_layout = ppt.slide_layouts[1]
slide = ppt.slides.add_slide(slide_layout)
# 采用相同的方式插入文字
title = slide.shapes.title
title.text = "我是标题2"
subtitle = slide.placeholders[1]
subtitle.text = "正文框2"

# 保存
ppt.save("write_text.pptx")



# 追加文字
ppt = Presentation("write_text.pptx")

# 第一页幻灯片
slide0 = ppt.slides[0]
# 获取第一页幻灯片所有的占位符
placeholder = slide0.shapes.placeholders
# 在第二个占位符对象中添加新的段落
new_paragraph = placeholder[1].text_frame.add_paragraph()
# 追加新文字
new_paragraph.text = "追加的新文字"
ppt.save("write_text2.pptx")
