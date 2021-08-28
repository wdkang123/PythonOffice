from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE


ppt = Presentation()
# 空白布局
layout = ppt.slide_layouts[6]
# 添加空白布局的幻灯片
slide = ppt.slides.add_slide(layout)
# 预设位置以及大小
left = Inches(5)
top = Inches(5)
width = Inches(5)
height = Inches(5)

# left top 为相对位置
# width height 为文本框大小
textbox = slide.shapes.add_textbox(left, top, width, height)
textbox.text = "这是一个新的文本框"
# 添加新的段落
new_paragraph = textbox.text_frame.add_paragraph()
new_paragraph.text = "文本框中第二段内容"


# 向ppt中添加图片
left = Inches(0)
top = Inches(0)
width = Inches(3)
height = Inches(3)
img_path = "1.jpg"
# 插入
pic = slide.shapes.add_picture(img_path, left, top, width, height)


# 向ppt中添加形状
left = Inches(1)
top = Inches(2)
width = Inches(1.8)
height = Inches(1)
# 插入形状
shape = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
# 在形状中加入文字
shape.text = "第一步"

for i in range(2, 6):
  # 移动位置
  left = left + width - Inches(0.3)
  # 插入形状
  shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
  shape.text = f"第{i}步"


ppt.save("add_new_content.pptx")
